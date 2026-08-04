#!/usr/bin/env python3
"""Generate the structure benchmark: ordinary Unicode documents with labelled contents.

Why this exists
---------------
Every number this repository publishes is measured on **legacy-encoded** documents. That
is the moat, and it is the right thing to measure — but it says nothing about what viparse
does with an ordinary Unicode `.docx` or PDF, which is most of what anyone will feed it.
Twenty minutes of unmeasured spot-checking on 2026-08-04 turned up three defects, so
"probably fine" was not a defensible answer.

What makes this different from the transcript corpus
----------------------------------------------------
The headline accuracy figure is circular and this file says so twice: the transcripts and
the conversion tables were derived from the same documents, so it measures
self-consistency as much as correctness.

This benchmark is not circular, and the reason is structural rather than a matter of care.
It does not compare against a transcript at all. It plants **labels** — numbered
paragraphs, named headings, a table with a known header — and asks whether they come back
in the right order, at the right level, and attached to the right things. The generator
shares no code with the parser, and a label like ``Đoạn số 07`` cannot be talked into
looking correct. Nothing here can be improved by editing the ground truth, because the
ground truth is a counting argument.

Its real weakness is the opposite one: **these documents are generated, not found.** A
generator emits clean, well-formed files, so this measures the parser against the easy
half of the world. It is a floor, not a ceiling — a defect it finds is real, a defect it
misses proves nothing.

    pip install python-docx openpyxl python-pptx reportlab
    python3 scripts/make_structure_bench.py --out structure/documents/
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

# A Unicode-capable TTF is required for the PDFs: reportlab's built-in Type 1 fonts have
# no Vietnamese repertoire, and would silently emit black boxes for every diacritic.
FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
    "C:/Windows/Fonts/arial.ttf",
]

VI = "Tổng sản phẩm trong nước ước tính tăng 6,42% so với cùng kỳ năm trước, khu vực dịch vụ đóng góp 48,7% vào mức tăng chung."
EN = "Figures are preliminary and subject to revision; seasonal adjustment follows the X-13ARIMA-SEATS procedure and rates are annualised."

SECTIONS = [
    ("Tình hình kinh tế vĩ mô", VI),
    ("Đầu tư và xuất khẩu", VI),
    ("Methodology and caveats", EN),
]

# Long enough that it cannot fit in one chunk at any sensible budget. A four-row table
# always fits, so a benchmark built on one would report "table: ok" while the defect it
# is looking for — data rows arriving without their header — never gets a chance to occur.
_INDICATORS = [
    ("Tăng trưởng GDP", "5,66%", "6,42%", "+0,76"),
    ("Xuất khẩu (tỷ USD)", "171,2", "195,4", "+14,2%"),
    ("Nhập khẩu (tỷ USD)", "158,9", "180,1", "+13,3%"),
    ("Lạm phát bình quân", "3,77%", "3,54%", "−0,23"),
    ("Thất nghiệp thành thị", "2,24%", "2,18%", "−0,06"),
    ("Vốn FDI đăng ký", "9,27", "11,84", "+27,7%"),
    ("Vốn FDI thực hiện", "6,28", "7,15", "+13,9%"),
    ("Khách quốc tế (triệu)", "4,63", "5,12", "+10,6%"),
    ("Bán lẻ hàng hoá", "8,15%", "8,84%", "+0,69"),
    ("Sản xuất công nghiệp", "5,91%", "7,54%", "+1,63"),
    ("Nông nghiệp", "3,34%", "3,68%", "+0,34"),
    ("Thu ngân sách", "39,6%", "51,2%", "+11,6"),
]

TABLE = [["Chỉ tiêu", "Quý I", "Quý II", "Thay đổi"], *[list(r) for r in _INDICATORS]]


def paragraph(index: int, body: str) -> str:
    """A paragraph that announces its own position, so order is checkable by counting."""
    return f"Đoạn số {index:02d}. {body}"


def _font() -> str:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for path in FONT_CANDIDATES:
        if Path(path).exists():
            pdfmetrics.registerFont(TTFont("Uni", path))
            return "Uni"
    print(
        "no Unicode TTF found; PDFs would lose every diacritic. Looked in:\n  "
        + "\n  ".join(FONT_CANDIDATES),
        file=sys.stderr,
    )
    raise SystemExit(2)


def build_docx(path: Path, count: int) -> dict:
    import docx

    document = docx.Document()
    document.add_heading("Báo cáo tình hình kinh tế quý II năm 2026", level=1)
    index = 1
    for title, body in SECTIONS:
        document.add_heading(title, level=2)
        for _ in range(count // len(SECTIONS)):
            document.add_paragraph(paragraph(index, body))
            index += 1
    table = document.add_table(rows=0, cols=len(TABLE[0]))
    table.style = "Table Grid"
    for row in TABLE:
        for cell, value in zip(table.add_row().cells, row, strict=True):
            cell.text = value
    document.save(path)
    return {"paragraphs": index - 1, "headings": [s[0] for s in SECTIONS], "table": TABLE}


def build_xlsx(path: Path) -> dict:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Tổng hợp"
    for row in TABLE:
        sheet.append(row)
    notes = workbook.create_sheet("Ghi chú")
    notes.append([paragraph(1, EN)])
    workbook.save(path)
    return {"paragraphs": 1, "headings": ["Tổng hợp", "Ghi chú"], "table": TABLE}


def build_pptx(path: Path, count: int) -> dict:
    import pptx
    from pptx.util import Inches

    presentation = pptx.Presentation()
    index = 1
    for title, body in SECTIONS:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.text = paragraph(index, body)
        index += 1
        for _ in range(count // len(SECTIONS) - 1):
            frame.add_paragraph().text = paragraph(index, body)
            index += 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Số liệu"
    shape = slide.shapes.add_table(
        len(TABLE), len(TABLE[0]), Inches(0.5), Inches(1.8), Inches(9), Inches(3)
    )
    for r, row in enumerate(TABLE):
        for c, value in enumerate(row):
            shape.table.cell(r, c).text = value
    presentation.save(path)
    return {"paragraphs": index - 1, "headings": [s[0] for s in SECTIONS], "table": TABLE}


def build_pdf(path: Path, count: int, columns: int) -> dict:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    from reportlab.platypus.doctemplate import BaseDocTemplate, PageTemplate
    from reportlab.platypus.frames import Frame

    font = _font()
    styles = getSampleStyleSheet()
    for name in ("Title", "Heading1", "Heading2", "BodyText"):
        styles[name].fontName = font

    story: list[object] = [Paragraph("Báo cáo tình hình kinh tế quý II năm 2026", styles["Title"])]
    index = 1
    for title, body in SECTIONS:
        story.append(Paragraph(title, styles["Heading2"]))
        for _ in range(count // len(SECTIONS)):
            story.append(Paragraph(paragraph(index, body), styles["BodyText"]))
            index += 1
    story += [
        Spacer(1, 10),
        Table(
            TABLE,
            style=TableStyle(
                [("GRID", (0, 0), (-1, -1), 0.5, colors.black), ("FONTNAME", (0, 0), (-1, -1), font)]
            ),
        ),
    ]

    if columns == 1:
        SimpleDocTemplate(str(path), pagesize=A4).build(story)
    else:
        width, height = A4
        margin, gap = 40, 18
        column_width = (width - 2 * margin - gap * (columns - 1)) / columns
        frames = [
            Frame(margin + i * (column_width + gap), margin, column_width, height - 80, id=str(i))
            for i in range(columns)
        ]
        document = BaseDocTemplate(str(path), pagesize=A4)
        document.addPageTemplates([PageTemplate(id="cols", frames=frames)])
        document.build(story)
    return {"paragraphs": index - 1, "headings": [s[0] for s in SECTIONS], "table": TABLE}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path, required=True)
    ap.add_argument(
        "--paragraphs",
        type=int,
        default=36,
        help="labelled paragraphs per document; must overflow one PDF column for the "
        "multi-column cases to test anything (default 36)",
    )
    args = ap.parse_args()
    args.out.mkdir(parents=True, exist_ok=True)
    n = args.paragraphs

    manifest = {
        "structured.docx": build_docx(args.out / "structured.docx", n),
        "structured.xlsx": build_xlsx(args.out / "structured.xlsx"),
        "structured.pptx": build_pptx(args.out / "structured.pptx", n),
        "one_column.pdf": build_pdf(args.out / "one_column.pdf", n, columns=1),
        "two_column.pdf": build_pdf(args.out / "two_column.pdf", n, columns=2),
        "three_column.pdf": build_pdf(args.out / "three_column.pdf", n, columns=3),
    }
    (args.out / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    for name, spec in manifest.items():
        print(f"  {name:22} {spec['paragraphs']} paragraphs, {len(spec['headings'])} headings")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
